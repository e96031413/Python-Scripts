# REF https://stackoverflow.com/questions/55497789/find-a-word-in-multiple-powerpoint-files-python/55763992#55763992
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import os

path = "./"

files = [x for x in os.listdir(path) if x.endswith(".pptx")] 


def CheckRecursivelyForText(shpthissetofshapes):
    for shape in shpthissetofshapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            checkrecursivelyfortext(shape.shapes)
        else:
            if hasattr(shape, "text"):
                shape.text = shape.text.lower()
                if "what_ever_you_want_to_find" in shape.text:
                    print(eachfile)
                    print("----------------------")
                else :
                    print("No text found in these PPTs")
                    print("----------------------")
                    break

for eachfile in files:
    prs = Presentation(path + eachfile) 
    for slide in prs.slides:
        CheckRecursivelyForText(slide.shapes)